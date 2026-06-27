# -*- coding: utf-8 -*-
"""
포에시아 회사소개서 v3.00 생성 스크립트 (2026-06-26)
- 원본(v2.02) 보존, 신규 파일 출력: 포에시아_회사소개서_v3.00_260626.pptx
- 본문 카피: 회사소개서_v3.00_본문_260626.md 기준
- 디자인: 16:9, Malgun Gothic, 포에시아 레드 테마. 기존 images/ 자산 재활용.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
IMG = os.path.join(HERE, "images")

# 색상
RED   = RGBColor(0xC8, 0x10, 0x2E)
DARK  = RGBColor(0x2B, 0x2B, 0x35)
GRAY  = RGBColor(0x5A, 0x5A, 0x60)
LGRAY = RGBColor(0x88, 0x88, 0x90)
LIGHT = RGBColor(0xF4, 0xF4, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PINK  = RGBColor(0xFB, 0xE9, 0xEC)
FONT  = "Malgun Gothic"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)
ML = Inches(0.9)               # 좌측 마진
CW = Inches(13.333 - 1.8)      # 컨텐츠 폭

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, l, t, w, h, fill=None, line=None):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def _style(run, size, bold, color, font=FONT):
    run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = color; run.font.name = font
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {}); rPr.append(el)
        el.set("typeface", font)


def tb(s, l, t, w, h, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT):
    box = s.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.paragraphs[0].alignment = align
    return tf


def para(tf, text, size=16, bold=False, color=GRAY, align=PP_ALIGN.LEFT,
         first=False, space_after=4, bullet=False, line=1.1):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.line_spacing = line
    pre = "•  " if bullet else ""
    run = p.add_run(); run.text = pre + text
    _style(run, size, bold, color)
    return p


def footer(s):
    f = tb(s, ML, Inches(7.02), CW, Inches(0.35))
    para(f, "AI 기술로 일상을 혁신합니다.  ·  POESIA", 9, False, LGRAY, first=True)


def kicker_title(s, kicker, title, lead=None, dark=False):
    tcol = WHITE if dark else DARK
    k = tb(s, ML, Inches(0.55), CW, Inches(0.4))
    para(k, kicker, 12, True, RED, first=True)
    t = tb(s, ML, Inches(0.92), CW, Inches(0.8))
    para(t, title, 30, True, tcol, first=True)
    rect(s, ML, Inches(1.72), Inches(0.7), Inches(0.06), fill=RED)
    if lead:
        lo = tb(s, ML, Inches(1.85), CW, Inches(0.5))
        para(lo, lead, 15, False, LGRAY, first=True)


def cards3(s, top, items, height=Inches(3.7)):
    gap = Inches(0.35)
    cw = Emu(int((CW - gap * 2) / 3))
    for i, it in enumerate(items):
        l = Emu(int(ML + i * (cw + gap)))
        fill = it.get("fill", LIGHT)
        rect(s, l, top, cw, height, fill=fill)
        inner = tb(s, Emu(int(l + Inches(0.28))), Emu(int(top + Inches(0.28))),
                   Emu(int(cw - Inches(0.56))), Emu(int(height - Inches(0.56))))
        if it.get("badge"):
            para(inner, it["badge"], 11, True, RED, first=True, space_after=6)
            para(inner, it["title"], 19, True, DARK, space_after=8)
        else:
            para(inner, it["title"], 19, True, DARK, first=True, space_after=8)
        if it.get("desc"):
            para(inner, it["desc"], 13, False, GRAY, space_after=10, line=1.25)
        for b in it.get("bullets", []):
            para(inner, b, 12.5, False, RED if False else GRAY, bullet=True, space_after=3, line=1.15)


# ───────────────────────── S1 표지 ─────────────────────────
s = slide()
rect(s, 0, 0, EMU_W, EMU_H, fill=RED)
t = tb(s, ML, Inches(2.0), CW, Inches(2.2))
para(t, "POESIA", 60, True, WHITE, first=True, space_after=2)
para(t, "Company Profile", 36, True, WHITE, space_after=18)
para(t, "SI · 서비스 · AI   |   AI 기술로 일상을 혁신합니다", 18, False, WHITE, space_after=2)
para(t, "You First You Always", 14, False, RGBColor(0xF2, 0xC8, 0xCE))
b = tb(s, ML, Inches(6.5), CW, Inches(0.8))
para(b, "경기도 남양주시 화도읍 마석중앙로37번길 45, 504호-D152호(별나라프라자)", 11, False, WHITE, first=True, space_after=2)
para(b, "E-mail  hello@smilewithpoesia.com   ·   www.smilewithpoesia.com", 11, False, WHITE)

# ───────────────────────── S2 Overview ─────────────────────────
s = slide()
kicker_title(s, "OVERVIEW", "회사 개요",
             "검증된 수행 경험의 SI 파트너이자, 직접 만들고 운영하는 서비스 회사")
rows = [
    ("회사명", "주식회사 포에시아 (POESIA)"),
    ("대표이사", "조연우"),
    ("설립일", "2023. 03. 02"),
    ("사업자등록번호", "475-87-02542"),
    ("소재지", "경기도 남양주시 화도읍 마석중앙로37번길 45, 504호-D152호(별나라프라자)"),
    ("연락처", "TEL [확인]  ·  hello@smilewithpoesia.com  ·  www.smilewithpoesia.com"),
    ("사업분야", "정보통신업 / 전문·과학·기술서비스업 (응용·시스템 SW 개발·공급, 학술·연구용역)"),
]
tbl_top = Inches(2.55)
tw = Inches(8.4)
tbl = s.shapes.add_table(len(rows), 2, ML, tbl_top, tw, Inches(3.5)).table
tbl.columns[0].width = Inches(2.2); tbl.columns[1].width = Inches(6.2)
tbl.first_row = False; tbl.horz_banding = False
for r, (k, v) in enumerate(rows):
    for c, txt in enumerate((k, v)):
        cell = tbl.cell(r, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = PINK if c == 0 else WHITE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.12); cell.margin_top = Inches(0.04); cell.margin_bottom = Inches(0.04)
        tf = cell.text_frame; tf.word_wrap = True
        para(tf, txt, 12.5, c == 0, DARK if c == 0 else GRAY, first=True)
# 정체성 박스(우측) — 신용등급 제외(결산·신용평가 미갱신)
bx = rect(s, Inches(9.55), tbl_top, Inches(2.9), Inches(3.5), fill=DARK)
bi = tb(s, Inches(9.75), Inches(2.75), Inches(2.5), Inches(3.1), anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
para(bi, "기업부설연구소", 20, True, RGBColor(0xFF, 0x5B, 0x7F), first=True, align=PP_ALIGN.CENTER, space_after=2)
para(bi, "2024.08 설립 · 자체 R&D 역량", 11, False, RGBColor(0xCF, 0xCF, 0xD6), align=PP_ALIGN.CENTER, space_after=16)
para(bi, "AI를 SI·서비스에 결합하는", 12, True, WHITE, align=PP_ALIGN.CENTER, space_after=1)
para(bi, "통합 · 운영 · 보안 역량", 12, True, WHITE, align=PP_ALIGN.CENTER)
footer(s)

# ───────────────────────── S3 Why 포에시아 ─────────────────────────
s = slide()
kicker_title(s, "WHY POESIA", "우리는 '파는 회사'가 아니라 '실행하고 운영하는 회사'입니다")
cards3(s, Inches(2.3), [
    {"badge": "01 DELIVERY", "title": "실행력", "fill": LIGHT,
     "desc": "공공·금융·BPO·플랫폼 등 실제 프로젝트를 수행해 온 경험.",
     "bullets": ["요구분석 → 설계 → 개발 → 운영", "끝까지 책임지는 수행"]},
    {"badge": "02 OPERATION", "title": "운영 경험", "fill": PINK,
     "desc": "제안서가 아니라 실제 서비스를 직접 만들고 운영합니다.",
     "bullets": ["자체 서비스 AXDevHub", "뮤직 플랫폼 구축·운영", "제품·운영 노하우 내재화"]},
    {"badge": "03 CAPABILITY", "title": "AI 결합", "fill": LIGHT,
     "desc": "AI 상향평준화 시대, 진짜 경쟁력은 모델 보유가 아닙니다.",
     "bullets": ["통합·운영·도메인·보안 역량", "SI·서비스에 AI를 '결합'"]},
], height=Inches(3.9))
footer(s)

# ───────────────────────── S4 사업영역 ─────────────────────────
s = slide()
kicker_title(s, "BUSINESS", "사업영역 — 3대 사업축",
             "실행하는 SI, 직접 운영하는 서비스, 이를 받치는 AI 역량")
cards3(s, Inches(2.35), [
    {"badge": "01 · SI · SM (매출 기반)", "title": "SI · 시스템 구축", "fill": LIGHT,
     "desc": "공공·금융·유통·제조 정보시스템의 구축·고도화·유지운영.",
     "bullets": ["시스템 통합(SI)·유지운영(SM)", "AICC·챗봇·데이터 융복합", "요구분석→설계→개발→운영"]},
    {"badge": "02 · SERVICE (성장)", "title": "자체 서비스", "fill": PINK,
     "desc": "직접 만들고 운영하는 서비스로 성장합니다.",
     "bullets": ["AXDevHub — 개발 생태계 허브", "구독형 AI 서비스(AICC)", "운영 기반 부가가치"]},
    {"badge": "03 · SOLUTION (차별화)", "title": "AI 역량 · 기술 자산", "fill": LIGHT,
     "desc": "SI·서비스에 바로 결합하는 AI 기술 자산.",
     "bullets": ["음성(STT·TTS)·보안(Defence)", "생성형 챗봇·콜봇(RAG/LLM)", "AICC·업무 자동화"]},
], height=Inches(3.95))
footer(s)

# ───────────────────────── S5 주요 수행실적 ─────────────────────────
s = slide()
kicker_title(s, "TRACK RECORD", "주요 수행실적",
             "제안이 아니라, 실제로 만든 결과로 말합니다")
data = [
    ("사례", "고객 / 구분", "역할"),
    ("뮤직 플랫폼 구축 및 운영", "음악 서비스 플랫폼", "구축 및 운영"),
    ("상담지원시스템 개발·운영", "글로벌 BPO", "개발 및 운영 참여"),
    ("상담 자동평가(AutoQA) 개발", "카드사", "상담 후처리 자동화"),
    ("민원신고 챗봇 개발", "공공기관", "개발 참여"),
    ("혜택정보 제공 콜봇", "카드사 연구과제", "MRC·RAG 기반 개발"),
    ("위치기반 서비스 개발", "세스코", "개발 참여"),
]
t5 = s.shapes.add_table(len(data), 3, ML, Inches(2.5), CW, Inches(3.7)).table
t5.columns[0].width = Inches(5.0); t5.columns[1].width = Inches(3.5); t5.columns[2].width = Inches(3.03)
t5.first_row = False; t5.horz_banding = False
for r, row in enumerate(data):
    head = (r == 0)
    for c, txt in enumerate(row):
        cell = t5.cell(r, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RED if head else (LIGHT if r % 2 else WHITE)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.14)
        tf = cell.text_frame; tf.word_wrap = True
        para(tf, txt, 13 if not head else 13, head or c == 0,
             WHITE if head else (DARK if c == 0 else GRAY), first=True)
nb = tb(s, ML, Inches(6.35), CW, Inches(0.4))
para(nb, "※ 고객사 실명·성과 수치 공개 범위는 확인 후 확정. 역할 표기는 실제 수행 범위 기준.", 10, False, LGRAY, first=True)
footer(s)

# ───────────────────────── S6 AXDevHub (강조 밴드) ─────────────────────────
s = slide()
rect(s, 0, 0, EMU_W, EMU_H, fill=RED)
k = tb(s, ML, Inches(0.7), CW, Inches(0.4))
para(k, "OUR SERVICE", 13, True, RGBColor(0xF6, 0xC8, 0xCE), first=True)
t = tb(s, ML, Inches(1.05), CW, Inches(0.9))
para(t, "AXDevHub", 40, True, WHITE, first=True)
d = tb(s, ML, Inches(1.95), Inches(11.2), Inches(1.1))
para(d, "포에시아가 직접 만들고 운영하는 SI·SW 개발 생태계 허브.", 17, True, WHITE, first=True, space_after=4)
para(d, "사업·과제부터 프로젝트·인력·커뮤니티까지 한곳에서 — 개발자와 기업을 잇습니다.", 15, False, RGBColor(0xF7, 0xDD, 0xE1))
feats = ["사업·과제 인텔리전스", "프로젝트 매칭", "인재풀", "개발자 커뮤니티"]
gap = Inches(0.3); fw = Emu(int((CW - gap * 3) / 4))
for i, f in enumerate(feats):
    l = Emu(int(ML + i * (fw + gap)))
    rect(s, l, Inches(3.35), fw, Inches(1.0), fill=RGBColor(0xD8, 0x3A, 0x53))
    fi = tb(s, l, Inches(3.35), fw, Inches(1.0), anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    para(fi, f, 14, True, WHITE, first=True, align=PP_ALIGN.CENTER)
mean = tb(s, ML, Inches(4.7), CW, Inches(1.1))
para(mean, "사업적 의미", 14, True, WHITE, first=True, space_after=4)
para(mean, "· 포에시아의 B2C 시장 진입 거점", 14, False, RGBColor(0xF7, 0xDD, 0xE1), space_after=3)
para(mean, "· 직접 운영을 통한 제품·운영 노하우 내재화 → SI·영업 자산으로 환류", 14, False, RGBColor(0xF7, 0xDD, 0xE1))
cta = rect(s, ML, Inches(6.1), Inches(3.3), Inches(0.7), fill=WHITE)
ci = tb(s, ML, Inches(6.1), Inches(3.3), Inches(0.7), anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
para(ci, "AXDevHub 바로가기  ↗", 15, True, RED, first=True, align=PP_ALIGN.CENTER)

# ───────────────────────── S7 구독형 AI 서비스 ─────────────────────────
s = slide()
kicker_title(s, "SERVICE", "구독형 AI 서비스 (AICC)",
             "초기 투자 부담 없이, 필요한 만큼 구독형으로")
cards3(s, Inches(2.35), [
    {"title": "구독형 STT/TA", "fill": LIGHT,
     "desc": "실시간·배치 음성인식과 자연어 처리. 유형분류·감정분석·요약으로 상담 후처리를 지원.",
     "bullets": ["실시간/배치 음성인식", "유형분류·감정분석·요약", "콜센터·금융·의료·공공"]},
    {"title": "구독형 상담지원", "fill": PINK,
     "desc": "실시간 상담 음성인식, 관리자 모니터링, KMS 연동·생성형 상담 가이드.",
     "bullets": ["실시간 상담 음성인식", "KMS 연동 상담 가이드", "생성형 상담 가이드"]},
    {"title": "구독형 콜봇", "fill": LIGHT,
     "desc": "보이스게이트웨이·STT·TTS·생성형 답변 기반 24시간 무중단 AI 음성 상담.",
     "bullets": ["내부문서 기반 AI 답변(RAG)", "할루시네이션 방지·상담사 연동", "금융·헬스케어·이커머스·공공"]},
], height=Inches(4.0))
footer(s)

# ───────────────────────── S8 AI 역량 한눈에 ─────────────────────────
s = slide()
kicker_title(s, "CAPABILITY", "AI 기술 역량 한눈에",
             "단일 제품이 아니라, 4계층을 조합해 고객 도메인에 맞게 구축하는 통합 역량")
layers = [
    ("기반제품 (Engine & Model)", "PoB-STT · PoB-TTS · PoB-Vision · PoB-Defence · PoB-LLM · PoB-VoiceGW-PBX · PoB-Record"),
    ("기술제품 (Solution)", "PoS-Callbot · PoS-ChatBot-Hybrid · PoS-Assist · PoS-AutoQA · PoS-STT/TA · PoS-Messenger · PoS-RPA"),
    ("응용제품 (Application)", "PoA-AICC · PoA-AI Agent · PoA-HyperAutomation"),
    ("플랫폼 (Platform)", "PoP-MLOps · PoP-AutoML · PoP-DevOps"),
]
yt = Inches(2.45); rh = Inches(0.95); gap = Inches(0.18)
for i, (head, items) in enumerate(layers):
    y = Emu(int(yt + i * (rh + gap)))
    rect(s, ML, y, Inches(3.2), rh, fill=RED)
    hi = tb(s, ML, y, Inches(3.2), rh, anchor=MSO_ANCHOR.MIDDLE)
    para(hi, head, 14, True, WHITE, first=True, align=PP_ALIGN.CENTER)
    rect(s, Inches(4.2), y, Inches(8.23), rh, fill=LIGHT)
    bi = tb(s, Inches(4.4), y, Inches(7.9), rh, anchor=MSO_ANCHOR.MIDDLE)
    para(bi, items, 13.5, False, DARK, first=True, line=1.15)
footer(s)

# ───────────────────────── S9 플래그십① AICC·콜봇 ─────────────────────────
s = slide()
kicker_title(s, "FLAGSHIP ①", "AICC · 콜봇 — 생성형·RAG",
             "단순 안내를 넘어, 생성형 AI 상담으로")
left = tb(s, ML, Inches(2.45), Inches(7.4), Inches(4.2))
para(left, "PoA-AICC — 통합 채널 AI 콜센터", 16, True, RED, first=True, space_after=4)
para(left, "전화·챗봇·콜봇 채널 통합 + 감정분석 기반 Voice Bot + 상담지원·생성형 상담 가이드.", 13.5, False, GRAY, space_after=2, line=1.25)
para(left, "구성: STT · TTS · Defence · VoiceGW · ChatBot · Assist · Callbot", 12, False, LGRAY, space_after=14)
para(left, "PoS-Callbot — 전문상담 콜봇", 16, True, RED, space_after=4)
for b in ["내부문서·매뉴얼 기반 생성형 답변(RAG), 할루시네이션 방지",
          "24시간 무중단 자동 상담 → 상담사 업무 절감",
          "기존 콜인프라 활용, 소규모도 합리적 비용 구축",
          "상담사 연동·부재 시 녹음/아웃바운드 지원"]:
    para(left, b, 13, False, GRAY, bullet=True, space_after=3, line=1.2)
para(left, "적용: 금융·헬스케어·이커머스·공공·IT/SaaS", 12, False, LGRAY, space_after=2)
img = os.path.join(IMG, "product_callbot.png")
if os.path.exists(img):
    s.shapes.add_picture(img, Inches(9.5), Inches(2.9), width=Inches(3.0))
footer(s)

# ───────────────────────── S10 플래그십② 보안·신분증 ─────────────────────────
s = slide()
kicker_title(s, "FLAGSHIP ②", "보안 · 신분증인식 — 방어 가능한 차별점",
             "AI 시대일수록 더 중요한 '안전한 데이터 활용'")
left = tb(s, ML, Inches(2.45), Inches(7.4), Inches(4.2))
para(left, "PoB-Defence — 민감정보 보안", 16, True, RED, first=True, space_after=4)
para(left, "정형·비정형 데이터 내 민감정보(주민·계좌·카드·여권·면허·주소·전화) 탐지 → 단/양방향 암호화·마스킹. 욕설·개인정보 필터, REST API 연동.", 13, False, GRAY, space_after=14, line=1.25)
para(left, "PoB-Vision — 신분증인식(자체 OCR)", 16, True, RED, space_after=4)
para(left, "신분증 4종(주민·면허·외국인·여권) 자동 분류·인식, 개인정보 자동 마스킹.", 13, False, GRAY, space_after=3, line=1.2)
para(left, "인식률 95% 이상 (주민 99.31% · 면허 98.67% · 외국인 100% · 여권 99.68%)", 13, True, DARK, space_after=2)
img = os.path.join(IMG, "product_infoguard.png")
if os.path.exists(img):
    s.shapes.add_picture(img, Inches(9.5), Inches(2.9), width=Inches(3.0))
# 성능 박스
rect(s, ML, Inches(6.2), CW, Inches(0.55), fill=DARK)
pb = tb(s, ML, Inches(6.2), CW, Inches(0.55), anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
para(pb, "성능 요약   ·   STT 95%+ (42,500시간 학습)   ·   TTS MOS 4.3+   ·   신분증 95%+",
     13, True, WHITE, first=True, align=PP_ALIGN.CENTER)

# (S11 신뢰지표 슬라이드 제거 — 결산·신용평가 미갱신으로 신용등급/신뢰지표 비표기.
#  제품 성능 수치는 S10 하단 '성능 요약' 박스로 대체, 기업부설연구소는 S2 표기.)

# ───────────────────────── S12 마무리 ─────────────────────────
s = slide()
rect(s, 0, 0, EMU_W, EMU_H, fill=RED)
t = tb(s, ML, Inches(2.3), CW, Inches(2.0))
para(t, "Complete The Value", 40, True, WHITE, first=True, space_after=6)
para(t, "포에시아와 함께 가치 있는 비즈니스를 만드세요", 24, True, WHITE)
b = tb(s, ML, Inches(6.0), CW, Inches(1.2))
para(b, "주식회사 포에시아 · 대표 조연우 · 사업자등록번호 475-87-02542", 12, False, WHITE, first=True, space_after=2)
para(b, "경기도 남양주시 화도읍 마석중앙로37번길 45, 504호-D152호(별나라프라자)", 12, False, WHITE, space_after=2)
para(b, "TEL [확인]  ·  FAX [확인]  ·  hello@smilewithpoesia.com  ·  www.smilewithpoesia.com", 12, False, WHITE)

out = os.path.join(HERE, "포에시아_회사소개서_v3.00_260626.pptx")
prs.save(out)
print("saved:", out, "| slides:", len(prs.slides._sldIdLst))
